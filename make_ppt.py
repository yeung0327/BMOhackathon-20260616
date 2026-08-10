from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_COLOR = RGBColor(0x0a, 0x0a, 0x1a)
CYAN = RGBColor(0x00, 0xd4, 0xff)
GOLD = RGBColor(0xf5, 0x9e, 0x0b)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0xa0, 0xae, 0xc0)
LIGHT_BLUE = RGBColor(0x66, 0xcc, 0xff)


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_multi_text(slide, left, top, width, height, lines):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(8)
    return tf


# ===== P1: 封面 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
add_multi_text(slide, Inches(1.5), Inches(2), Inches(10), Inches(4), [
    ("🌱 根与芽 · 知识宇宙", 44, GOLD, True, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("文档关联 · 探索发现 · 对话溯源 —— 让知识自己长出连接", 22, LIGHT_BLUE, False, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("百度 BMO 黑客松 2026", 16, GRAY, False, PP_ALIGN.CENTER),
    ("大模型运营部 | 单人 + AI 协作", 14, GRAY, False, PP_ALIGN.CENTER),
])

# ===== P2: 痛点 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
         "❓ 问题：文档孤岛", 32, CYAN, True)
add_multi_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), [
    ("每个新人入职，面对的是飞书上几十篇分散的文档——", 22, WHITE, False, PP_ALIGN.LEFT),
    ("活动复盘、产品手册、渠道资料、会议纪要……", 20, GRAY, False, PP_ALIGN.LEFT),
    ("", 14, WHITE, False, PP_ALIGN.LEFT),
    ("❌  文档之间没有连接", 22, WHITE, False, PP_ALIGN.LEFT),
    ("❌  找信息只能靠关键词碰运气", 22, WHITE, False, PP_ALIGN.LEFT),
    ("❌  跨文档的关联完全靠人脑记忆", 22, WHITE, False, PP_ALIGN.LEFT),
    ("", 14, WHITE, False, PP_ALIGN.LEFT),
    (""我们不缺文档，缺的是文档之间的连接。"", 26, GOLD, True, PP_ALIGN.CENTER),
])

# ===== P3: 方案 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
         "💡 方案：上传文档 → AI自动生成知识网络", 30, CYAN, True)
add_multi_text(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(5), [
    ("三大核心能力", 24, WHITE, True, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("🔍  探索式浏览", 22, LIGHT_BLUE, True, PP_ALIGN.LEFT),
    ("    双击节点，关联知识自动展开", 18, GRAY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("💬  对话式检索", 22, LIGHT_BLUE, True, PP_ALIGN.LEFT),
    ("    AI回答 + 图谱实体实时高亮溯源", 18, GRAY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("📄  一键摘要", 22, LIGHT_BLUE, True, PP_ALIGN.LEFT),
    ("    点击文档节点，30秒生成摘要", 18, GRAY, False, PP_ALIGN.LEFT),
])
add_multi_text(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(5), [
    ("技术架构", 24, WHITE, True, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("前端  React + TypeScript + Neo4j NVL", 16, GRAY, False, PP_ALIGN.LEFT),
    ("后端  Python + FastAPI", 16, GRAY, False, PP_ALIGN.LEFT),
    ("图数据库  Neo4j AuraDB", 16, GRAY, False, PP_ALIGN.LEFT),
    ("LLM  DeepSeek（模型无关架构）", 16, GRAY, False, PP_ALIGN.LEFT),
    ("嵌入  all-MiniLM-L6-v2（本地）", 16, GRAY, False, PP_ALIGN.LEFT),
    ("部署  Docker Compose", 16, GRAY, False, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("⚡ 换模型只需改一个环境变量", 16, GOLD, False, PP_ALIGN.LEFT),
])

# ===== P4: Live Demo =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_multi_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(4), [
    ("🎬  Live Demo", 48, CYAN, True, PP_ALIGN.CENTER),
    ("", 16, WHITE, False, PP_ALIGN.CENTER),
    ("457 个节点  ·  1,030 条关系  ·  零人工标注", 24, GOLD, False, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("（切换到浏览器展示）", 18, GRAY, False, PP_ALIGN.CENTER),
])

# ===== P5: 业务价值 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
         "📈 业务价值：真实数据，真实场景", 30, CYAN, True)
add_multi_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(5.5), [
    ("用的是部门真实运营文档，不是Wikipedia，不是模拟数据", 20, GOLD, True, PP_ALIGN.LEFT),
    ("", 12, WHITE, False, PP_ALIGN.LEFT),
    ("👤  新人入职        30分钟建立业务全貌（原来需要1-2周）", 20, WHITE, False, PP_ALIGN.LEFT),
    ("📋  活动策划        一眼看到历史渠道与效果关联", 20, WHITE, False, PP_ALIGN.LEFT),
    ("📝  复盘沉淀        上传即入网，自动被后人发现", 20, WHITE, False, PP_ALIGN.LEFT),
    ("🔄  知识激活        文档越多，网络越密，价值复利增长", 20, WHITE, False, PP_ALIGN.LEFT),
    ("", 14, WHITE, False, PP_ALIGN.LEFT),
    ("与传统方案的区别：", 22, WHITE, True, PP_ALIGN.LEFT),
    ("飞书搜索 → 你得知道关键词 | 根与芽 → 点击即发现", 18, GRAY, False, PP_ALIGN.LEFT),
    ("纯AI问答 → 等你提问 | 根与芽 → 让你浏览", 18, GRAY, False, PP_ALIGN.LEFT),
    ("向量相似 → 匿名相似度 | 根与芽 → 有名字的关系", 18, GRAY, False, PP_ALIGN.LEFT),
])

# ===== P6: 技术亮点 + 数据 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
         "🔧 技术亮点", 30, CYAN, True)
add_multi_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), [
    ("① 模型无关架构 —— 一个环境变量切换 DeepSeek / 文心 / 任意模型", 20, WHITE, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("② 图谱质量优化 —— 1587碎片 → 457有效节点（-71%），零孤立节点", 20, WHITE, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("③ 渐进式探索 —— 双击展开两跳邻居，不过载用户", 20, WHITE, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("④ 对话-图谱联动 —— AI回答中的实体在图上实时高亮放大", 20, WHITE, False, PP_ALIGN.LEFT),
    ("", 14, WHITE, False, PP_ALIGN.LEFT),
    ("⏱️  5天   |   👤 1人   |   📄 6篇真实文档   |   🤖 AI辅助开发", 22, GOLD, True, PP_ALIGN.CENTER),
])

# ===== P7: 收尾 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5), [
    ("落地路径", 28, WHITE, True, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("接入飞书文档同步  +  切换文心模型  +  部署内网", 22, GRAY, False, PP_ALIGN.CENTER),
    ("预计两周可上线", 22, LIGHT_BLUE, True, PP_ALIGN.CENTER),
    ("", 24, WHITE, False, PP_ALIGN.CENTER),
    ("", 24, WHITE, False, PP_ALIGN.CENTER),
    (""我们不缺文档，缺的是文档之间的连接。"", 28, GOLD, True, PP_ALIGN.CENTER),
    ("根与芽，让知识自己长出网络。", 24, CYAN, False, PP_ALIGN.CENTER),
    ("", 20, WHITE, False, PP_ALIGN.CENTER),
    ("谢谢！", 32, WHITE, True, PP_ALIGN.CENTER),
])

output_path = "/Users/yangqianqian/Desktop/BMOhackathon-20260616/路演PPT-根与芽知识宇宙.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
